# from pptx import Presentation
#
# # Укажите путь к файлу презентации
# file_path = '/Users/Stanislav_Egorov/Documents/GitHub/blitz-croco-words/src/croco-blitz-source/Osennyaya_igra_3.pptx'
#
# # Открываем файл презентации
# presentation = Presentation(file_path)
#
# # Теперь вы можете работать с объектом presentation
# # Например, выводим количество слайдов
# # print(f'Количество слайдов в презентации: {len(presentation.slides)}')
#
# # Вы можете итерироваться по слайдам
# # for slide in presentation.slides:
# #     print(f'Слайд {slide.slide_id}')
#
# prs = Presentation('Osennyaya_igra_3.pptx')

# for slide in prs.slides:
#     for shape in slide.shapes:
#         if not shape.has_text_frame:
#             continue
#         print(shape.text)


# def extract_words_without_spaces(file_path):
#     prs = Presentation(file_path)
#
#     words_without_spaces = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, 'text'):
#                 text = shape.text.strip()
#                 words = text.split()
#                 if len(words) == 1:
#                     words_without_spaces.append(words[0])
#
#     return words_without_spaces
#
#
# extracted_words_without_spaces = extract_words_without_spaces(file_path)
# for word in extracted_words_without_spaces:
#     print(word)

import zipfile
from pathlib import Path
from pptx import Presentation
from random import choice
from pyaspeller import YandexSpeller

# Константы
SRC = 'src'

with zipfile.ZipFile('/Users/Stanislav_Egorov/Documents/GitHub/blitz-croco-words/src/croco-blitz-source.zip',
                     'r') as zip_file:
    PPTX_FILES_LIST: list[str] = [
        "Osennyaya_igra_3.pptx",
        "Zimnyaya_igra_1.pptx",
        "Osennyaya_igra_12.pptx",
        "Osennyaya_igra_11.pptx",
        "Osennyaya_igra_10.pptx",
        "Osennyaya_igra_9.pptx",
        "Osennyaya_igra_6.pptx",
        "Osennyaya_igra_5.pptx",
        "Osennyaya_igra_4.pptx"
    ]


def open_pptx_with_pptx_library(presentation_path):
    prs = Presentation(presentation_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                print(shape.text)


def extract_words_from_pptx(presentation_path):
    words = []
    prs = Presentation(presentation_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text
                if ' ' in text or ':' in text or 'БЛИЦ-КРОКОДИЛ' in text:
                    continue
                words.append(text.strip())
    return words


def write_txt(filename: str, words: list[str]):
    try:
        with open(filename, 'w', encoding='utf-8') as f:
            for word in words:
                f.write(word + '\n')
        print(f"Слова успешно записаны в файл {filename}")
    except IOError as e:
        print(f"Ошибка при записи в файл: {e}")


# def only_one_word(line: str) -> bool:
#     return ' ' not in line

# Задача 2024.10.30.02
#
# Напишите тест, который проверяет, что в строке находится всего одно слово.
def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text


# def test_only_one_word():
#     assert only_one_word(' ') == False
#     assert only_one_word('word') == True


def main() -> None:
    # Извлекаем файл презентации
    # extract_pptx_from_zip(SRC, FILENAME_ZIP, EXTRACTION_DIR)

    # Открываем и выводим содержимое выбранной презентации
    presentation_path = '/Users/Stanislav_Egorov/Documents/GitHub/blitz-croco-words/src/croco-blitz-source'
    # open_pptx_with_pptx_library(str(presentation_path))

    # Выбираем случайный файл презентации
    pptx_file_name = choice(PPTX_FILES_LIST)
    presentation_path = Path(presentation_path) / pptx_file_name
    # open_pptx_with_pptx_library(str(presentation_path))

    # Извлекаем слова из презентации
    extracted_words = extract_words_from_pptx(str(presentation_path))

    # Проверка на наличие только одного слова в строке


    # Записываем извлеченные слова в txt-файл
    write_txt('words.txt', extracted_words)


main()
