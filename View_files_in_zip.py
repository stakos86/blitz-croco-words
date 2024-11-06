import os
import zipfile

from pptx import Presentation


def extract_pptx_from_zip(zip_filename, pptx_filename, extract_path):
    # Распаковываем zip-файл
    with zipfile.ZipFile(zip_filename, 'r') as zip_ref:
        zip_ref.extractall(extract_path)

    # Служебное сообщение для проверки, что файлы распакованы
    print(f"Файлы из {zip_filename} распакованы в {extract_path}")

    # Путь к PPTX файлу
    full_pptx_path = os.path.join(extract_path, pptx_filename)

    # Проверка, существует ли PPTX файл
    if not os.path.isfile(full_pptx_path):
        print(f"Файл {full_pptx_path} не найден!")
        return

    # Загружаем PPTX презентацию и выводим текст слайдов
    presentation = Presentation(full_pptx_path)

    for slide_number, slide in enumerate(presentation.slides):
        print(f"Слайд {slide_number + 1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)


# Пример использования
zip_filename = '/Users/Stanislav_Egorov/Documents/GitHub/blitz-croco-words/src/croco-blitz-source.zip'  # Имя zip-файла
pptx_filename = 'Osennyaya_igra_3.pptx'  # Имя PPTX файла внутри zip
extract_path = '/Users/Stanislav_Egorov/Documents/GitHub/blitz-croco-words/src/croco-blitz-source'  # Путь для извлечения файлов

extract_pptx_from_zip(zip_filename, pptx_filename, extract_path)
